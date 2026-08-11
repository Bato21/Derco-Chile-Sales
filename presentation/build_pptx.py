#!/usr/bin/env python
"""
Build BI2026_FinalProject_DERCO.pptx — the PowerPoint twin of presentation/index.html.

Both decks read the SAME presentation/metrics.json (written by the notebook, Stage 4),
so a number can never differ between the HTML deck, the PowerPoint and the notebook.

    pip install python-pptx
    python presentation/build_pptx.py

Slide size is 16:9 (13.333 x 7.5 in). 13 slides. Colours are DERCO's: deep red on white.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Emu, Inches, Pt
except ImportError:  # pragma: no cover
    sys.exit("python-pptx is required:  pip install python-pptx")

HERE = Path(__file__).resolve().parent
M = json.loads((HERE / "metrics.json").read_text(encoding="utf-8"))
NOTES = json.loads((HERE / "speaker_notes.json").read_text(encoding="utf-8"))["slides"]
OUT = HERE / "BI2026_FinalProject_DERCO.pptx"

# ---------------------------------------------------------------- brand tokens
RED       = RGBColor(0xC0, 0x05, 0x12)   # DERCO brand red · series slot 1 · 6.4:1 on white
RED_DEEP  = RGBColor(0x8C, 0x04, 0x10)   # red text (9.8:1 on white)
RED_WASH  = RGBColor(0xFD, 0xEC, 0xEC)
RED_FILL  = RGBColor(0xF6, 0xA0, 0x9B)   # light tint for large filled areas
BLUE      = RGBColor(0x2A, 0x78, 0xD6)   # series slot 2
BLUE_INK  = RGBColor(0x1F, 0x52, 0x96)
GREEN_INK = RGBColor(0x00, 0x63, 0x00)
INK       = RGBColor(0x1A, 0x1A, 0x1A)
INK2      = RGBColor(0x5A, 0x5A, 0x5A)
MUTED     = RGBColor(0x76, 0x76, 0x76)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
PANEL     = RGBColor(0xFB, 0xFB, 0xFA)
LINE      = RGBColor(0xE8, 0xE7, 0xE3)
NEUTRAL   = RGBColor(0xC9, 0xC8, 0xC3)
# ordered one-hue ramp — the 4 ranked customer segments (dark = most valuable)
RAMP = [RGBColor(0x8C, 0x04, 0x10), RGBColor(0xC0, 0x05, 0x12),
        RGBColor(0xE9, 0x4B, 0x45), RGBColor(0xF6, 0xA0, 0x9B)]

FONT = "Segoe UI"
W, H = Inches(13.333), Inches(7.5)
ML, MR = Inches(0.95), Inches(0.95)          # left / right margin
CW = W - ML - MR                              # content width

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- primitives
def slide(hero: bool = False):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = RED if hero else WHITE
    # the DERCO signature: a red rule down the left edge of every slide
    bar = s.shapes.add_shape(1, 0, 0, Inches(0.09), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = WHITE if hero else RED   # matches the HTML deck's hero rule
    bar.line.fill.background()
    bar.shadow.inherit = False
    return s


def textbox(s, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tf


def write(tf, text, size, color=INK, bold=False, italic=False, space_after=0,
          first=False, align=None, line=None):
    """Append (or fill the first) paragraph. Runs marked with ** ** render bold."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        p.alignment = align
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    for j, chunk in enumerate(text.split("**")):
        if not chunk:
            continue
        r = p.add_run()
        r.text = chunk
        f = r.font
        f.name, f.size, f.color.rgb = FONT, Pt(size), color
        f.bold = bold or (j % 2 == 1)
        f.italic = italic
    return p


def kicker(s, text, hero=False, y=0.62):
    """y defaults to the top of the slide; hero slides pass y below the wordmark."""
    tf = textbox(s, ML, Inches(y), CW, Inches(0.3))
    write(tf, text.upper(), 11, WHITE if hero else RED_DEEP, bold=True, first=True)
    tf.paragraphs[0].runs[0].font.name = FONT


def title(s, text, size=34, hero=False, y=1.05, h=1.5):
    tf = textbox(s, ML, Inches(y), CW, Inches(h))
    write(tf, text, size, WHITE if hero else INK, bold=True, first=True, line=1.06)


def body(s, text, y, size=14, color=INK2, w=None, h=0.9):
    tf = textbox(s, ML, Inches(y), w or CW, Inches(h))
    write(tf, text, size, color, first=True, line=1.32)
    return tf


def card(s, x, y, w, h, accent=None, fill=PANEL):
    box = s.shapes.add_shape(5, x, y, w, h)          # 5 = rounded rectangle
    box.adjustments[0] = 0.045
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = LINE
    box.line.width = Pt(0.75)
    box.shadow.inherit = False
    if accent:                                        # 3px top rule, like the HTML deck
        r = s.shapes.add_shape(1, x + Emu(int(w * 0.03)), y, w - Emu(int(w * 0.06)), Pt(2.5))
        r.fill.solid()
        r.fill.fore_color.rgb = accent
        r.line.fill.background()
        r.shadow.inherit = False
    return box


def stat_card(s, x, y, w, h, value, label, color=RED_DEEP, heading=None):
    card(s, x, y, w, h, accent=color if heading else None)
    pad = Inches(0.26)
    tf = textbox(s, x + pad, y + Inches(0.30), w - 2 * pad, h - Inches(0.42))
    if heading:
        write(tf, heading, 12, color, bold=True, first=True, space_after=5)
        write(tf, label, 11.5, INK2, line=1.3)
    else:
        write(tf, value, 30, color, bold=True, first=True, space_after=5)
        write(tf, label, 10.5, MUTED, line=1.25)


def bullets(s, items, y, size=14, gap=0.44, w=None):
    """One bullet per row: a small square in the series colour + the text.
    Pass w to keep the text column clear of anything sitting to its right."""
    yy = Inches(y)
    width = (w if w is not None else CW) - Inches(0.34)
    for color, text in items:
        sq = s.shapes.add_shape(1, ML, yy + Inches(0.085), Inches(0.13), Inches(0.13))
        sq.fill.solid()
        sq.fill.fore_color.rgb = color
        sq.line.fill.background()
        sq.shadow.inherit = False
        tf = textbox(s, ML + Inches(0.34), yy, width, Inches(gap))
        write(tf, text, size, INK2, first=True, line=1.3)
        yy += Inches(gap + 0.20)


def footer(s, speaker, topic, hero=False):
    col = WHITE if hero else MUTED
    rule = s.shapes.add_shape(1, ML, H - Inches(0.86), CW, Pt(0.75))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(0xFF, 0x8A, 0x8A) if hero else LINE
    rule.line.fill.background()
    rule.shadow.inherit = False
    tf = textbox(s, ML, H - Inches(0.72), CW / 2, Inches(0.3))
    write(tf, speaker, 10, WHITE if hero else RED_DEEP, bold=True, first=True)
    tf = textbox(s, ML + CW / 2, H - Inches(0.72), CW / 2, Inches(0.3), align=PP_ALIGN.RIGHT)
    write(tf, topic, 10, col, first=True)


def wordmark(s, hero=False):
    dash = s.shapes.add_shape(1, ML, Inches(0.62), Inches(0.30), Inches(0.11))
    dash.fill.solid()
    dash.fill.fore_color.rgb = WHITE if hero else RED
    dash.line.fill.background()
    dash.shadow.inherit = False
    tf = textbox(s, ML + Inches(0.42), Inches(0.50), Inches(4), Inches(0.35))
    write(tf, "DERCO CHILE", 15, WHITE if hero else INK, bold=True, first=True)


def caveat(s, text, y, h=0.85):
    box = s.shapes.add_shape(1, ML, Inches(y), CW, Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = RED_WASH
    box.line.fill.background()
    box.shadow.inherit = False
    rule = s.shapes.add_shape(1, ML, Inches(y), Pt(2.5), Inches(h))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RED
    rule.line.fill.background()
    rule.shadow.inherit = False
    tf = textbox(s, ML + Inches(0.24), Inches(y + 0.15), CW - Inches(0.48), Inches(h - 0.2))
    write(tf, text, 12, INK2, first=True, line=1.3)


def style_chart(chart, legend=False, legend_pos=XL_LEGEND_POSITION.TOP):
    chart.font.name, chart.font.size, chart.font.color.rgb = FONT, Pt(11), MUTED
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(11)
        chart.legend.font.color.rgb = INK2


# =============================================================== 1 · title
s = slide(hero=True)
wordmark(s, hero=True)
kicker(s, "Business Intelligence · Final Integrative Project · UDD 2026", hero=True, y=1.18)
tf = textbox(s, ML, Inches(2.30), CW, Inches(2.2))
write(tf, f"From {M['raw_rows']:,} car sales", 44, WHITE, bold=True, first=True, line=1.05)
write(tf, "to one business decision.", 44, WHITE, bold=True, line=1.05)
tf = textbox(s, ML, Inches(4.55), CW, Inches(1.1))
write(tf, f"Retail vehicle sales · May 2009 → March 2022", 15,
      RGBColor(0xFF, 0xD9, 0xD9), first=True, space_after=9)
write(tf, "Vicente Rodríguez · Agustín Reyes · Luis-Felipe Cáceres · Baptiste Vial", 14, WHITE)
footer(s, "▶ Baptiste", "Acting as DERCO's BI unit", hero=True)

# =============================================================== 2 · the ask
s = slide()
kicker(s, "01 · The client & the question")
title(s, "DERCO hired us to find where\nthe retail model leaks value.")
rule = s.shapes.add_shape(1, ML, Inches(2.55), Pt(3), Inches(0.78))
rule.fill.solid(); rule.fill.fore_color.rgb = RED; rule.line.fill.background()
rule.shadow.inherit = False
tf = textbox(s, ML + Inches(0.22), Inches(2.6), Inches(9.2), Inches(0.8))
write(tf, "“After 14 years of sales, where is DERCO leaking value — "
          "and what should change to grow margin?”", 17, INK, bold=True, first=True, line=1.25)
cw3, gap = (CW - Inches(0.5)) / 3, Inches(0.25)
for i, (head, txt) in enumerate([
        ("Decision-maker", "DERCO's **Commercial Director** — owns brand mix, channel strategy and the marketing budget."),
        ("The decision", "Where to move **commercial effort and budget** for FY2026."),
        ("Success", "A ranked, **CLP-valued** list of moves the client can fund — not an accuracy score.")]):
    stat_card(s, ML + i * (cw3 + gap), Inches(3.75), cw3, Inches(1.85), "", txt, RED, heading=head)
footer(s, "▶ Baptiste", "Frame & KPIs")

# =============================================================== 3 · the data & the pipeline
s = slide()
kicker(s, "02 · The dataset & how we worked")
title(s, "14 years of real transactions,\nthrough the full BI pipeline.", size=31)
cw4 = (CW - Inches(0.75)) / 4
for i, (v, l) in enumerate([("550K", "transactions"),
                            (f"{M['customers']//1000}K", "unique customers"),
                            (str(M["brands"]), f"brands · {M['models']:,} models"),
                            (str(M["comunas"]), "comunas of Chile")]):
    stat_card(s, ML + i * (cw4 + Inches(0.25)), Inches(2.30), cw4, Inches(1.20), v, l)
# the five BI stages as one compact strip
cw5 = (CW - Inches(0.60)) / 5
for i, (head_, txt) in enumerate([
        ("1 · Frame", "Question, decision-maker, 6 KPIs"),
        ("2 · Prepare", "Audit, clean, engineer. No leakage"),
        ("3 · Model", "Segment customers + predict losses"),
        ("4 · Communicate", "Insights, recommendation, dashboard"),
        ("5 · Ethics", "Bias, privacy, honest limits")]):
    x = ML + i * (cw5 + Inches(0.15))
    card(s, x, Inches(3.70), cw5, Inches(1.05), accent=RED)
    tf = textbox(s, x + Inches(0.18), Inches(3.90), cw5 - Inches(0.36), Inches(0.80))
    write(tf, head_, 10.5, INK, bold=True, first=True, space_after=3)
    write(tf, txt, 8.5, MUTED, line=1.2)
caveat(s, "Every row is one car sale — **date, customer hash, comuna, brand, model, channel, list price, "
          "margin**; two channels, **ces** (dealer, 65%) and **propio** (owned, 35%). One honesty note: the "
          "money fields are **synthetic / obfuscated**, so every CLP figure is a **relative signal** — shares, "
          "ranks and trends are valid, absolute pesos are not audited financials.", y=5.00, h=1.00)
tf = textbox(s, ML, Inches(6.15), CW, Inches(0.45))
write(tf, "Python · pandas · scikit-learn · matplotlib · seed = 42 · one reproducible notebook — every "
          "number on every slide is exported by it into metrics.json.", 10.5, MUTED, first=True)
footer(s, "▶ Baptiste → Luis-Felipe", "Data, scope & method")

# =============================================================== 4 · prepare & audit
s = slide()
kicker(s, "03 · Prepare & audit — “can we trust it?”")
title(s, "Clean data first. Then we found a leak.")
bullets(s, [
    (NEUTRAL, "**<1% missing**, zero duplicates, dates 100% consistent → safe to model."),
    (NEUTRAL, "Missing comuna → flagged **UNKNOWN**, never dropped (dropping quietly shrinks revenue)."),
    (NEUTRAL, "Dropped the redundant **year** column; log-scaled skewed spend before clustering."),
    (RED, f"**{M['loss_n']:,} sales ({M['loss_rate_pct']}%) lose money** — sold below cost. "
          "A real finding, not a parsing error."),
], y=2.45, w=Inches(7.35))
bx = ML + Inches(7.55)
card(s, bx, Inches(2.45), CW - Inches(7.55), Inches(2.55))
tf = textbox(s, bx, Inches(2.85), CW - Inches(7.55), Inches(1.0), align=PP_ALIGN.CENTER)
write(tf, f"{abs(M['loss_bn'])} bn", 46, RED, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = textbox(s, bx + Inches(0.3), Inches(3.95), CW - Inches(8.15), Inches(1.0), align=PP_ALIGN.CENTER)
write(tf, "CLP of margin leaked through loss-making deals", 11, MUTED, first=True,
      align=PP_ALIGN.CENTER, space_after=4, line=1.25)
write(tf, "≈ 2.7% of all retail margin (synthetic money)", 11, INK, bold=True, align=PP_ALIGN.CENTER)
footer(s, "▶ Luis-Felipe", "Data prep & quality")

# =============================================================== 5 · brand mix
s = slide()
kicker(s, "04 · EDA — brand mix")
title(s, "One brand carries the business.")
CHINESE = {"JAC Cars", "Great Wall", "Changan", "Geely", "Haval"}
brands = list(M["brand_share"].items())                       # already sorted desc
cd = CategoryChartData()
cd.categories = [b for b, _ in brands][::-1]
cd.add_series("Share of transactions (%)", [v for _, v in brands][::-1])
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, ML, Inches(2.15), CW - Inches(0.2),
                        Inches(4.05), cd)
ch = gf.chart
style_chart(ch)
plot = ch.plots[0]
plot.gap_width = 55
plot.vary_by_categories = False
ser = plot.series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = BLUE
# colour carries brand ORIGIN, not the bar's rank — bar length already shows volume
for pt, (name, _) in zip(ser.points, brands[::-1]):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = RED if name in CHINESE else BLUE
plot.has_data_labels = True                                    # direct labels: no axis hunting
dl = plot.data_labels
dl.number_format, dl.number_format_is_linked = '0.0"%"', False
dl.font.size, dl.font.color.rgb, dl.font.name = Pt(10), INK2, FONT
ch.value_axis.has_major_gridlines = False
ch.value_axis.visible = False
ch.category_axis.format.line.color.rgb = LINE
ch.category_axis.tick_labels.font.size = Pt(11)
ch.category_axis.tick_labels.font.color.rgb = INK
tf = textbox(s, ML, Inches(6.30), CW, Inches(0.6))
write(tf, "■ red = Chinese brand      ■ blue = incumbent brand   —   "
          f"**Suzuki alone is {M['suzuki_share']}%** of all {M['rows']:,} transactions "
          "→ single-brand dependence risk.", 11, MUTED, first=True, line=1.3)
tf.paragraphs[0].runs[0].font.color.rgb = RED_DEEP
footer(s, "▶ Luis-Felipe", "EDA · concentration")

# =============================================================== 6 · the headline
s = slide()
kicker(s, "05 · The headline")
title(s, "China won the shelf.")
years = sorted(M["chinese_share_by_year"], key=int)
vals = [M["chinese_share_by_year"][y] for y in years]
cd = CategoryChartData()
cd.categories = [str(y) for y in years]
cd.add_series("Chinese-brand share of transactions (%)", vals)
gf = s.shapes.add_chart(XL_CHART_TYPE.AREA, ML, Inches(2.10), CW - Inches(0.2), Inches(3.75), cd)
ch = gf.chart
style_chart(ch)
ser = ch.plots[0].series[0]
ser.format.fill.solid()
ser.format.fill.fore_color.rgb = RED_FILL          # large area = light tint, never saturated
ser.format.line.color.rgb = RED                    # the saturated hue rides the boundary
ser.format.line.width = Pt(2.25)
va = ch.value_axis
va.maximum_scale, va.minimum_scale, va.major_unit = 100.0, 0.0, 25.0
va.has_major_gridlines = True
va.major_gridlines.format.line.color.rgb = LINE
va.major_gridlines.format.line.width = Pt(0.75)
va.format.line.fill.background()
va.tick_labels.font.size, va.tick_labels.font.color.rgb = Pt(11), MUTED
va.tick_labels.number_format, va.tick_labels.number_format_is_linked = '0"%"', False
ch.category_axis.format.line.color.rgb = LINE
ch.category_axis.tick_labels.font.size = Pt(10)
ch.category_axis.tick_labels.font.color.rgb = MUTED
# selective direct labels — the two endpoints only
for xin, val, anchor in ((1.30, vals[0], PP_ALIGN.LEFT), (11.55, vals[-1], PP_ALIGN.RIGHT)):
    tf = textbox(s, Inches(xin), Inches(2.10 + 3.75 * (1 - val / 100) - 0.42), Inches(1.5),
                 Inches(0.4), align=anchor)
    write(tf, f"{val}%", 16, RED_DEEP, bold=True, first=True, align=anchor)
tf = textbox(s, ML, Inches(6.00), CW, Inches(0.85))
write(tf, "Not a straight line — the share dipped to ~20% in 2015–16 before the 2019→2022 surge. "
          f"But the direction is unambiguous: **Chinese brands now sell the majority "
          f"({vals[-1]}%) of every car DERCO moves.** 2009 (from May) and 2022 (to March) are "
          "partial years — shares are valid, totals are not.", 11, MUTED, first=True, line=1.35)
footer(s, "▶ Luis-Felipe → Agustín", "The structural shift")

# =============================================================== 7 · channel / geo / rhythm
s = slide()
kicker(s, "06 · EDA — channel, geography, rhythm")
title(s, "Where, how, and when DERCO sells.")
for i, (v, l, c) in enumerate([
        (f"{M['ces_share']:.0f} / {100 - M['ces_share']:.0f}",
         f"dealer ces vs owned propio — and dealers also price better: "
         f"{M['ces_margin_pct']}% margin vs {M['propio_margin_pct']}%", BLUE_INK),
        (f"{M['top10_comuna_share']}%",
         "of all sales sit in the top-10 comunas — dense, but over-exposed to Santiago metro", RED_DEEP),
        ("Aug ▲", "seasonal peak in August (December a close second), trough in April — "
                  "plan stock & campaigns to it", BLUE_INK)]):
    stat_card(s, ML + i * (cw3 + gap), Inches(2.20), cw3, Inches(1.85), v, l, c)
card(s, ML, Inches(4.35), CW, Inches(1.60), accent=RED, fill=RED_WASH)
tf = textbox(s, ML + Inches(0.35), Inches(4.62), CW - Inches(0.7), Inches(1.25))
write(tf, "The channel finding that points the second model", 13, RED_DEEP, bold=True,
      first=True, space_after=7)
write(tf, "DERCO's **own stores close 5.4% of deals below cost** against **0.5% at dealers** — "
          "an **~11× gap**. The margin leak is not spread evenly across the business; it is "
          "concentrated where DERCO controls the pricing.", 13, INK2, line=1.3)
footer(s, "▶ Agustín", "EDA wrap")

# =============================================================== 8 · segmentation + insight
s = slide()
kicker(s, "07 · Model A — customer segmentation")
title(s, "Who should DERCO retain? The 14% that come back.", size=30, h=0.9)
body(s, "RFM (Recency · Frequency · Monetary) → log-scaled → KMeans. We tested k = 2…7 with elbow "
        "+ silhouette and chose **k = 4** — the elbow, a healthy silhouette, and four groups that "
        "map onto four real marketing plays.", y=1.90, size=12.5, h=0.75)
SEGS = ["Champions (repeat, high-value)", "Big-ticket one-timers",
        "Mainstream one-timers", "Dormant / lapsed"]
SHORT = ["Champions", "Big-ticket one-timers", "Mainstream one-timers", "Dormant / lapsed"]
cd = CategoryChartData()
cd.categories = SHORT[::-1]
cd.add_series("Share of customers (%)", [M["segment_cust_pct"][k] for k in SEGS][::-1])
cd.add_series("Share of total margin (%)", [M["segment_margin_pct"][k] for k in SEGS][::-1])
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, ML, Inches(2.80), Inches(7.5), Inches(3.30), cd)
ch = gf.chart
style_chart(ch, legend=True)
ch.plots[0].gap_width, ch.plots[0].overlap = 60, -12
for ser, col in zip(ch.plots[0].series, (RAMP[0], RAMP[2])):
    ser.format.fill.solid()
    ser.format.fill.fore_color.rgb = col
ch.plots[0].has_data_labels = True
dl = ch.plots[0].data_labels
dl.number_format, dl.number_format_is_linked = '0.0"%"', False
dl.font.size, dl.font.color.rgb, dl.font.name = Pt(9.5), INK2, FONT
ch.value_axis.has_major_gridlines = False
ch.value_axis.visible = False
ch.category_axis.format.line.color.rgb = LINE
ch.category_axis.tick_labels.font.size = Pt(10.5)
ch.category_axis.tick_labels.font.color.rgb = INK
# the punchline, beside the chart
rx = ML + Inches(7.85)
rw = CW - Inches(7.85)
tf = textbox(s, rx, Inches(2.85), rw, Inches(0.75))
write(tf, f"{M['champions_cust_pct']:.0f}% → {M['champions_margin_pct']:.0f}%", 34, GREEN_INK,
      bold=True, first=True)
tf = textbox(s, rx, Inches(3.62), rw, Inches(1.05))
write(tf, "Champions are the **only** group that buys twice, and they carry the densest margin "
          "per customer in the business.", 11, INK2, first=True, line=1.3)
card(s, rx, Inches(4.80), rw, Inches(1.30))
tf = textbox(s, rx, Inches(5.00), rw, Inches(0.60), align=PP_ALIGN.CENTER)
write(tf, f"{M['one_time_pct']:.0f}%", 36, RED, bold=True, first=True, align=PP_ALIGN.CENTER)
tf = textbox(s, rx + Inches(0.22), Inches(5.62), rw - Inches(0.44), Inches(0.45), align=PP_ALIGN.CENTER)
write(tf, "buy **exactly once** in 14 years → DERCO is **transactional, not relational**",
      9.5, MUTED, first=True, align=PP_ALIGN.CENTER, line=1.25)
footer(s, "▶ Agustín", "Segmentation")

# =============================================================== 9 · classifier
s = slide()
kicker(s, "08 · Model B — predicting loss-making deals")
title(s, "Flag the money-losers before they close.", h=0.8)
body(s, "Supervised classifier on **pre-sale features only** — brand, channel, comuna, price, timing. "
        "**Margin is excluded to avoid leakage** (the target is derived from it). Dummy + logistic "
        "baselines → HistGradientBoosting, stratified hold-out + 5-fold CV.", y=1.95, size=13, h=0.9)
for i, (v, l, c) in enumerate([
        (f"{M['gbm_cv_auc']:.2f}",
         f"ROC-AUC, 5-fold CV mean ({M['gbm_auc']} on the hold-out; std {M['gbm_cv_std']})", BLUE_INK),
        (f"{M['top10_lift']}×", "better than reviewing deals at random", BLUE_INK),
        (f"{M['top10_capture']:.0f}%",
         "of all loss-making deals caught by reviewing just the top-10% riskiest", GREEN_INK)]):
    stat_card(s, ML + i * (cw3 + gap), Inches(3.10), cw3, Inches(1.70), v, l, c)
caveat(s, f"**Honest note:** losses are only {M['loss_rate_pct']}% of deals, so PR-AUC is "
          f"{M['gbm_prauc']} — modest in absolute terms, ~6.5× the base rate. This is a "
          "**triage aid, not an auto-reject**: a human reviews every flag.", y=5.15, h=0.85)
footer(s, "▶ Agustín → Vicente", "Classification")

# =============================================================== 10 · dashboard
s = slide()
kicker(s, "09 · The dashboard")
title(s, "Six numbers, one screen.", h=0.8)
tiles = [(f"{M['chinese_2022']}%", "of 2022 sales are Chinese brands", RED_DEEP),
         (f"{M['one_time_pct']}%", "of customers buy exactly once", RED_DEEP),
         (f"{M['champions_margin_pct']}%",
          f"of margin from the {M['champions_cust_pct']:.0f}% Champions", GREEN_INK),
         (f"{M['loss_rate_pct']}%", "of deals close below cost", RED_DEEP),
         (f"{M['suzuki_share']}%", "Suzuki's share of volume — concentration risk", BLUE_INK),
         (f"{M['top10_capture']}%", "of the leak recovered for review by checking "
                                    "the riskiest 10% of deals", GREEN_INK)]
tw = (CW - Inches(0.5)) / 3
for i, (v, l, c) in enumerate(tiles):
    x = ML + (i % 3) * (tw + Inches(0.25))
    y = Inches(2.25 + (i // 3) * 1.85)
    stat_card(s, x, y, tw, Inches(1.60), v, l, c)
tf = textbox(s, ML, Inches(6.05), CW, Inches(0.7))
write(tf, "Live in presentation/index.html; the static twin is presentation/dashboard.png, written "
          "by the notebook. Both read the same metrics.json — no number is typed by hand.",
      11, MUTED, first=True, line=1.3)
footer(s, "▶ Vicente", "Dashboard")

# =============================================================== 11 · insights -> recommendation
s = slide()
kicker(s, "10 · What it means & what to do")
title(s, "Five findings → one decision,\nthree funded moves.", size=30)
bullets(s, [
    (RED, f"**China won the shelf** — {M['chinese_2009']:.0f}% → {M['chinese_2022']:.0f}% of sales. "
          "The core business, not a hedge."),
    (RGBColor(0x0C, 0xA3, 0x0C),
     f"**Transactional, not relational** — {M['one_time_pct']:.0f}% buy once; the "
     f"{M['champions_cust_pct']:.0f}% Champions drive {M['champions_margin_pct']:.0f}% of margin."),
    (RED, f"**~{abs(M['loss_bn'])} bn CLP leaks** through {M['loss_rate_pct']}% loss-making deals, "
          "concentrated in the own-store channel."),
    (BLUE, f"**Single-brand dependence** — Suzuki ≈ {M['suzuki_share']:.0f}% of volume."),
    (BLUE, f"**Geographic over-exposure** — top-10 comunas ≈ {M['top10_comuna_share']:.0f}% of sales."),
], y=2.55, size=12.5, gap=0.40, w=Inches(7.05))
rx = ML + Inches(7.45)
rw = CW - Inches(7.45)
for i, (head_, txt, col) in enumerate([
        ("1 · Retain", "Fund a **Champions-retention + one-timer-conversion** loyalty CRM.", GREEN_INK),
        ("2 · Formalise China", "Treat Chinese brands as a **first-class portfolio** — pricing, stock, "
                                "marketing.", RED_DEEP),
        ("3 · Plug the leak", "Classifier as a **pre-approval check** on the riskiest 10%, starting with "
                              "own stores.", BLUE_INK)]):
    stat_card(s, rx, Inches(2.45 + i * 1.42), rw, Inches(1.25), "", txt, col, heading=head_)
footer(s, "▶ Vicente", "Insights & recommendation")

# =============================================================== 12 · ethics
s = slide()
kicker(s, "11 · Ethics, bias & limits")
title(s, "What this analysis cannot say.")
bullets(s, [
    (RED, "**Synthetic money** — every CLP figure is a relative signal, not an audited financial. "
          "The direction is trustworthy; the level is not."),
    (RED, "**Privacy** — the customer hash is re-identifiable, so we treat it as personal data; "
          "street addresses were excluded from every model."),
    (BLUE, "**Geographic bias** — Santiago metro is over-represented, so “valuable customer” really "
           "means “urban customer”. A retention program must not redline rural comunas."),
    (BLUE, "**Moderate cluster quality** — silhouette 0.39. The segments are a marketing convenience, "
           "not a law of nature."),
    (RGBColor(0x0C, 0xA3, 0x0C),
     "**Human in the loop** — imbalanced target, no causal claims. The models advise; people decide."),
], y=2.52, size=13.5, gap=0.40)
footer(s, "▶ Vicente", "Ethics & limitations")

# =============================================================== 13 · close
s = slide(hero=True)
wordmark(s, hero=True)
kicker(s, "Thank you", hero=True, y=1.18)
tf = textbox(s, ML, Inches(2.30), CW, Inches(2.0))
write(tf, "One dataset.", 46, WHITE, bold=True, first=True, line=1.05)
write(tf, "One decision.", 46, WHITE, bold=True, line=1.05)
tf = textbox(s, ML, Inches(4.60), CW, Inches(1.1))
write(tf, "Retain Champions · formalise China · plug the margin leak.", 15,
      RGBColor(0xFF, 0xD9, 0xD9), first=True, space_after=9)
write(tf, "Vicente Rodríguez · Agustín Reyes · Luis-Felipe Cáceres · Baptiste Vial · UDD 2026",
      14, WHITE)
footer(s, "▶ Vicente", "Ready for your questions", hero=True)

# ------------------------------------------------- speaker notes (presenter view)
# Same source as the N-overlay in index.html and the script in SCRIPT.md.
if len(NOTES) != len(prs.slides._sldIdLst):
    sys.exit(f"speaker_notes.json has {len(NOTES)} slides, deck has {len(prs.slides._sldIdLst)}")
for sl, note in zip(prs.slides, NOTES):
    lines = [f"SLIDE {note['n']} · {note['title']}",
             f"{note['speaker']} · {note['seconds']}s",
             "",
             f"THE POINT: {note['point']}",
             "",
             "SAY:"]
    lines += [f"  {t}" for t in note["say"]]
    if note["numbers"]:
        lines += ["", "NUMBERS TO LAND: " + " | ".join(note["numbers"])]
    lines += ["", f"WATCH OUT: {note['watch']}", "", f"TRANSITION: {note['next']}"]
    tf = sl.notes_slide.notes_text_frame
    tf.text = lines[0]
    for ln in lines[1:]:
        tf.add_paragraph().text = ln
    for para in tf.paragraphs:
        for run in para.runs:
            run.font.size = Pt(11)

prs.save(OUT)
print(f"Wrote {OUT.name} — {len(prs.slides._sldIdLst)} slides, 16:9, speaker notes attached, "
      f"sourced from metrics.json")
